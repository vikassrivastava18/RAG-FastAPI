import os, tempfile
from datetime import timedelta, datetime, timezone
from dotenv import load_dotenv

import jwt
from fastapi import HTTPException
from langchain_community.document_loaders import PyPDFLoader


# colour slide 
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from db.models import Book
from core.config import logger, password_hash

load_dotenv()
BASE_DIR = os.path.dirname(os.path.abspath(__file__))


def create_ppt_from_content(content: list[dict], textbook_name: str, book_id:int, db) -> str:
    try:
        book = db.query(Book).filter(Book.id == book_id).first()
        left_logo_path = os.path.join(BASE_DIR, book.logo1) if book.logo1 else None
        right_logo_path = os.path.join(BASE_DIR, book.logo2) if book.logo2 else None
        disclaimer = book.disclaimer if book.disclaimer else ""
        
        prs = Presentation()

        def set_background_color(slide, rgb_color: tuple[int, int, int]):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*rgb_color)

        def add_footer(slide):
            footer_top = Inches(6.8)
            footer_height = Inches(0.5)

            # Left logo
            try:
                slide.shapes.add_picture(left_logo_path, Inches(0.2), footer_top, height=footer_height)
            except Exception as e:
                logger.warning(f"Left logo error: {e}")

            # Right logo
            if right_logo_path:
                try:
                    slide.shapes.add_picture(right_logo_path, Inches(9.0), footer_top, height=footer_height)
                except Exception as e:
                    logger.warning(f"Right logo error: {e}")

            # Disclaimer
            left = Inches(1.5)
            width = Inches(7.0)
            height = Inches(0.6)
            textbox = slide.shapes.add_textbox(left, footer_top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for i, line in enumerate(disclaimer.splitlines()):
                para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                para.text = line.strip()
                para.font.size = Pt(9)
                para.font.color.rgb = RGBColor(100, 100, 100)
                para.alignment = 1  # Center

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = textbook_name
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""
        set_background_color(slide, (230, 240, 255))
        add_footer(slide)

        # Content Slides
        for slide_data in content:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            set_background_color(slide, (230, 240, 255))

            # Title box
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            title_tf = title_box.text_frame
            title_tf.text = slide_data['title']
            title_tf.paragraphs[0].font.size = Pt(24)
            title_tf.paragraphs[0].font.bold = True

            # Content box
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(4.8))
            tf = content_box.text_frame
            tf.word_wrap = True

            for bullet in slide_data['bullets']:
                para = tf.add_paragraph()
                para.text = bullet
                para.level = 0
                para.font.size = Pt(16)

            add_footer(slide)

        # Save the presentation
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name
    except Exception as e:
        print(f"Error {e}")
        raise Exception(f"Error in creating PPT: {e}")
    

def extract_text_from_pdf(book_name: str, pdf_folder:str='') -> str:

    """Extracts text from the first 10 pages of a given book's PDF file."""
    
    pdf_path = os.path.join(pdf_folder, f"{book_name}")
    MAX_PAGE= 15

    if not os.path.exists(pdf_path):
        raise HTTPException(status_code=404, detail=f"Book '{book_name}' not found in the database.")
    try:
        pdf_loader = PyPDFLoader(pdf_path)
        pages = pdf_loader.load()[:MAX_PAGE] 
        text = "\n".join([page.page_content for page in pages])
        return text

    except Exception as e:

        raise HTTPException(status_code=400, detail=f"Error extracting text from PDF: {str(e)}")
    

def verify_password(plain_password, hashed_password):
    return password_hash.verify(plain_password, hashed_password)


def create_access_token(data: dict, expires_delta: timedelta | None = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.now(timezone.utc) + expires_delta
    else:
        expire = datetime.now(timezone.utc) + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, os.getenv("SECRET_KEY"), algorithm=os.getenv("ALGORITHM"))
    return encoded_jwt
